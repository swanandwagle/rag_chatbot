
import PyPDF2
import docx
from pptx import Presentation
from pathlib import Path
from typing import Optional


class DocumentProcessor:
    """Service for processing and extracting text from documents"""
    
    @staticmethod
    async def extract_text(file_path: str) -> str:
        """
        Extract text from document based on file type
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Extracted text content
            
        Raises:
            ValueError: If file type is not supported
        """
        path = Path(file_path)
        extension = path.suffix.lower()
        
        if extension == ".pdf":
            return await DocumentProcessor._extract_from_pdf(file_path)
        elif extension in [".docx", ".doc"]:
            return await DocumentProcessor._extract_from_docx(file_path)
        elif extension in [".pptx", ".ppt"]:
            return await DocumentProcessor._extract_from_pptx(file_path)
        elif extension == ".txt":
            return await DocumentProcessor._extract_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")
    
    @staticmethod
    async def _extract_from_pdf(file_path: str) -> str:
        """Extract text from PDF file"""
        text = ""
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
        return text.strip()
    
    @staticmethod
    async def _extract_from_docx(file_path: str) -> str:
        """Extract text from DOCX file"""
        doc = docx.Document(file_path)
        text = "\n\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text])
        return text.strip()
    
    @staticmethod
    async def _extract_from_pptx(file_path: str) -> str:
        """Extract text from PPTX file"""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        return text.strip()
    
    @staticmethod
    async def _extract_from_txt(file_path: str) -> str:
        """Extract text from TXT file"""
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read().strip()

